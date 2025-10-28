import numpy as np, io, zipfile
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
from PIL import Image
import plotly.graph_objects as go

plt.rcParams.update({'figure.dpi':120, 'axes.grid':True})

def fig_time_series(t, y, labels, title):
    fig, ax = plt.subplots(figsize=(6,3.2))
    for i in range(y.shape[1]): ax.plot(t, y[:,i], label=labels[i])
    ax.set_xlabel('Time [s]'); ax.set_title(title); ax.legend(ncol=min(len(labels),3)); return fig
def fig_angle(t, ang):
    fig, ax = plt.subplots(figsize=(6,3.2)); ax.plot(t, ang*180/np.pi)
    ax.set_ylabel('Angle Error [deg]'); ax.set_xlabel('Time [s]'); ax.set_title('Attitude Error (Geodesic Angle)'); return fig
def fig_phase(w):
    fig, ax = plt.subplots(figsize=(3.2,3.2)); ax.plot(w[:,0], w[:,1]); ax.set_xlabel('ωx'); ax.set_ylabel('ωy'); ax.set_title('Phase Plane (ωx-ωy)'); return fig
def fig_energy(t,J,w):
    E = 0.5*np.sum(w*(J@w.T).T, axis=1); fig, ax = plt.subplots(figsize=(6,3.2)); ax.plot(t,E)
    ax.set_title('Rotational Kinetic Energy'); ax.set_xlabel('Time [s]'); ax.set_ylabel('E [J]'); return fig
def fig_pointing_sphere(P):
    u,v=np.mgrid[0:2*np.pi:40j,0:np.pi:20j]; xs=np.cos(u)*np.sin(v); ys=np.sin(u)*np.sin(v); zs=np.cos(v)
    frames=[go.Frame(data=[go.Scatter3d(x=[P[k,0]],y=[P[k,1]],z=[P[k,2]],mode='markers',marker=dict(size=4))], name=str(k)) for k in range(P.shape[0])]
    fig=go.Figure(data=[go.Surface(x=xs,y=ys,z=zs,opacity=0.15,showscale=False), go.Scatter3d(x=[P[0,0]],y=[P[0,1]],z=[P[0,2]],mode='markers')],
                  layout=go.Layout(scene=dict(aspectmode='cube'), title='Body +X Pointing on Unit Sphere',
                                   updatemenus=[dict(type='buttons',buttons=[dict(label='Play',method='animate',args=[None])])]))
    fig.frames=frames; return fig
def _collect_figs(t,q,w,tau,h,err,J):
    figs=[fig_time_series(t,w,['ωx','ωy','ωz'],'Body Rates'),
          fig_time_series(t,tau,['τx','τy','τz'],'Control Torque'),
          fig_time_series(t,h,['hx','hy','hz'],'Wheel Momentum'),
          fig_angle(t,err), fig_phase(w), fig_energy(t,J,w),
          fig_time_series(t,q[:,1:],['qx','qy','qz'],'Quaternion (vector part)'),
          fig_time_series(t,np.linalg.norm(w,axis=1)[:,None],['|ω|'],'Rate Magnitude'),
          fig_time_series(t,np.linalg.norm(tau,axis=1)[:,None],['|τ|'],'Torque Magnitude'),
          fig_time_series(t,np.linalg.norm(h,axis=1)[:,None],['|h_rw|'],'Wheel Momentum Magnitude')]
    derr=np.gradient(err,t); figs.append(fig_time_series(t,derr[:,None],['d(angle)/dt'],'Error Angle Rate'))
    effort=np.cumsum(np.abs(tau),axis=0)*(t[1]-t[0]); figs.append(fig_time_series(t,effort,['Ex','Ey','Ez'],'Cumulative Effort'))
    return figs
def make_exam_pack(t,q,w,tau,h,err,J,P):
    figs=_collect_figs(t,q,w,tau,h,err,J); mem=io.BytesIO(); z=zipfile.ZipFile(mem,'w',zipfile.ZIP_DEFLATED)
    for i,f in enumerate(figs,1): buf=io.BytesIO(); f.savefig(buf,format='png',bbox_inches='tight'); z.writestr(f'fig_{i:02d}.png',buf.getvalue()); plt.close(f)
    img_frames=[]; P2=P[:,:2]; step=max(1,P2.shape[0]//80)
    for k in range(0,P2.shape[0],step):
        fig,ax=plt.subplots(figsize=(3,3)); ax.plot(P2[:k+1,0],P2[:k+1,1]); ax.set_xlim(-1,1); ax.set_ylim(-1,1); ax.set_aspect('equal'); ax.set_title('Pointing Path (proj)')
        buf=io.BytesIO(); fig.savefig(buf,format='png',bbox_inches='tight'); plt.close(fig); from PIL import Image; img=Image.open(io.BytesIO(buf.getvalue())).convert('P',palette=Image.ADAPTIVE); img_frames.append(img)
    gif_buf=io.BytesIO()
    if len(img_frames)>1: img_frames[0].save(gif_buf,format='GIF',save_all=True,append_images=img_frames[1:],duration=60,loop=0); z.writestr('pointing.gif',gif_buf.getvalue())
    z.writestr('summary.md', "# DSAL Exam Pack Summary\n\n- Settling time (angle < 1 deg): computed in UI\n- Peak rate/torque: see figures\n- Energy decay indicates controller damping efficiency\n")
    z.close(); mem.seek(0); return mem.getvalue()
def make_exam_ppt(t,q,w,tau,h,err,J):
    from pptx import Presentation
    from pptx.util import Inches
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text="DeepSpace Attitude Lab — Results"; slide.placeholders[1].text="12 Figures • Attitude Control • Detumble / PID / LQR"
    figs=_collect_figs(t,q,w,tau,h,err,J)
    def fig_png(fig): import io; buf=io.BytesIO(); fig.savefig(buf,format='png',bbox_inches='tight',dpi=150); plt.close(fig); return buf.getvalue()
    img_w=4.5; left1=Inches(0.5); left2=Inches(5.0); top=Inches(1.0)
    for i in range(0,len(figs),2):
        s=prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.add_picture(io.BytesIO(fig_png(figs[i])), left1, top, width=Inches(img_w))
        if i+1<len(figs): s.shapes.add_picture(io.BytesIO(fig_png(figs[i+1])), left2, top, width=Inches(img_w))
    bio=io.BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()
